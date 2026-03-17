#!/usr/bin/env python3
"""Generate presentation slides (.pptx) for VC 162 final presentation.

Anthropic-inspired theme: warm cream backgrounds, clean typography, orange accents.
Import the .pptx into Google Slides for editing.

15-slide structure:
  1. Title
  2. Background & Motivation
  3. Connection to Readings (Abebe, Narayanan)
  4. Intervention & Voting Thresholds
  5. Datasets
  6. Pipeline
  7. Two Approaches: XGBoost (data-driven) vs Generative (theory-driven)
  8. XGBoost Method (detail)
  9. USA Map — County Approval Probability
  10. Generative Approval Model (done outside this project)
  11. Monte Carlo Simulation Engine
  12. Current Snapshot: What Vanishes? (generative)
  13. Projection: Future Growth Map (XGBoost)
  14. Projection: Future Growth Curves (XGBoost)
  15. Conclusion
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Anthropic-inspired palette ──────────────────────────────────────────────
BG = RGBColor(0xFA, 0xF7, 0xF2)
TITLE_CLR = RGBColor(0x19, 0x19, 0x2C)
BODY_CLR = RGBColor(0x2D, 0x2D, 0x3D)
MUTED_CLR = RGBColor(0x7A, 0x77, 0x70)
ACCENT = RGBColor(0xD4, 0x76, 0x3C)    # Warm orange
ACCENT2 = RGBColor(0x4A, 0x7C, 0x5E)   # Muted green
ACCENT3 = RGBColor(0x8B, 0x4A, 0x6B)   # Muted purple
LIGHT_BOX = RGBColor(0xED, 0xE8, 0xDF)
DARK_BOX = RGBColor(0x2D, 0x2D, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HDR = RGBColor(0x3D, 0x3D, 0x4D)
TABLE_ROW1 = RGBColor(0xF5, 0xF0, 0xE8)
TABLE_ROW2 = RGBColor(0xFA, 0xF7, 0xF2)

FONT = "Helvetica Neue"

W = Inches(13.333)
H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BODY_CLR, bold=False, alignment=PP_ALIGN.LEFT,
                line_spacing=1.2, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def rich_tf(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_para(tf, text, size=18, color=BODY_CLR, bold=False,
             align=PP_ALIGN.LEFT, after=6, before=0, italic=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    return p


def bullet(tf, text, size=16, color=BODY_CLR, bold=False, level=0, after=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.level = level
    p.space_after = Pt(after)
    return p


def box(slide, left, top, width, height, fill=LIGHT_BOX):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_card(slide, left, top, number, label, clr=ACCENT):
    box(slide, left, top, Inches(2.8), Inches(1.5), LIGHT_BOX)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                Inches(2.5), Inches(0.8), number,
                font_size=34, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.85),
                Inches(2.5), Inches(0.5), label,
                font_size=12, color=MUTED_CLR, alignment=PP_ALIGN.CENTER)


def img_placeholder(slide, left, top, width, height, label):
    box(slide, left, top, width, height, RGBColor(0xE5, 0xE0, 0xD8))
    add_textbox(slide, left + Inches(0.3), top + height / 2 - Inches(0.25),
                width - Inches(0.6), Inches(0.5), f"[Insert {label}]",
                font_size=14, color=MUTED_CLR, alignment=PP_ALIGN.CENTER)


# ── Slides ──────────────────────────────────────────────────────────────────

def slide_01(prs):
    """Title."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BOX)

    accent_line(s, Inches(1.5), Inches(1.0), Inches(2), ACCENT)
    add_textbox(s, Inches(1.5), Inches(1.3), Inches(10), Inches(2.0),
                "What If Tech Corporations Bore\nthe Cost of Consent?",
                font_size=42, color=WHITE, bold=True, line_spacing=1.15)
    add_textbox(s, Inches(1.5), Inches(3.3), Inches(9), Inches(0.8),
                "Modeling Data Center Growth Under Alternative Consent Regimes",
                font_size=20, color=MUTED_CLR)
    add_textbox(s, Inches(1.5), Inches(5.2), Inches(9), Inches(0.5),
                "Amrita  \u00b7  Lauren  \u00b7  Dhruv  \u00b7  Tyler",
                font_size=18, color=RGBColor(0xAA, 0xA7, 0xA0))
    add_textbox(s, Inches(1.5), Inches(5.8), Inches(9), Inches(0.5),
                "VC 162  \u00b7  Values in Computational Thinking  \u00b7  Spring 2026",
                font_size=14, color=MUTED_CLR)

    notes(s,
          "[Amrita]\n\n"
          "Our project asks: right now, if a community doesn\u2019t want a data center, "
          "they have to organize, hire lawyers, and sustain pressure campaigns to block "
          "projects already in motion. What if we flipped that \u2014 what if the company "
          "had to earn consent first?")


def slide_02(prs):
    """Background & Motivation."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Background & Motivation",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Stat cards
    stat_card(s, Inches(1.0), Inches(1.6), "10 \u2192 72 GW",
              "U.S. data center capacity\n2025 \u2192 2035", ACCENT)
    stat_card(s, Inches(4.2), Inches(1.6), "~600",
              "Facilities in Northern\nVirginia alone", ACCENT3)
    stat_card(s, Inches(7.4), Inches(1.6), "44%",
              "Americans who support DCs\nin their community", ACCENT2)

    _, tf = rich_tf(s, Inches(1.0), Inches(3.5), Inches(10.5), Inches(3.5))
    add_para(tf, "The asymmetry:", size=20, bold=True, color=TITLE_CLR, after=8)
    bullet(tf,
           "Data center siting is routine permitting \u2014 no formal consent "
           "requirement in most U.S. jurisdictions", size=16, after=6)
    bullet(tf,
           "Communities must reactively organize, fund legal challenges, "
           "and sustain pressure campaigns to block projects", size=16, after=6)
    bullet(tf,
           "JLARC (2024): benefits concentrate in a few counties while costs "
           "(electricity, water, land use) spread across regions", size=16, after=6)
    bullet(tf,
           "The burden of resistance falls entirely on communities, "
           "not on corporations proposing development", size=16,
           color=ACCENT, bold=True, after=6)

    notes(s,
          "[Amrita]\n\n"
          "Data centers are the physical backbone of AI. Industry forecasts project "
          "16\u201320 GW of new capacity per year. ABI Research projects U.S. capacity "
          "growing from 10.2 GW to 71.8 GW by 2035.\n\n"
          "But communities hosting them see rising electricity costs, aquifer depletion, "
          "and agricultural land conversion. Northern Virginia \u2014 home to the world's "
          "largest concentration \u2014 has roughly 600 facilities. Virginia's Joint "
          "Legislative Audit and Review Commission found that while data centers "
          "generate substantial tax revenue, the benefits and costs are distributed "
          "unevenly, and saturation effects are already visible in hyperconcentrated "
          "counties.\n\n"
          "The key point: in the current system, consent is passive. There's no formal "
          "threshold \u2014 communities have to reactively organize. The cost of saying "
          "'no' falls entirely on residents.")


def slide_03(prs):
    """Connection to Readings."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Connection to Readings",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Abebe box
    box(s, Inches(1.0), Inches(1.6), Inches(5.2), Inches(2.6), LIGHT_BOX)
    _, tf1 = rich_tf(s, Inches(1.3), Inches(1.75), Inches(4.7), Inches(2.3))
    add_para(tf1, "Abebe et al. (2020)", size=14, color=ACCENT, bold=True, after=4)
    add_para(tf1, "Roles for Computing in Social Change", size=20,
             color=TITLE_CLR, bold=True, after=10)
    add_para(tf1, "Formalizer:", size=15, bold=True, color=BODY_CLR, after=2)
    add_para(tf1, "Make the implicit power structure of data center "
             "siting explicit and measurable.", size=13, color=BODY_CLR, after=6)
    add_para(tf1, "Rebuttal:", size=15, bold=True, color=BODY_CLR, after=2)
    add_para(tf1, "Simulate alternatives that challenge the assumption "
             "the status quo is neutral.", size=13, color=BODY_CLR, after=0)

    # Narayanan box
    box(s, Inches(6.8), Inches(1.6), Inches(5.5), Inches(2.6), LIGHT_BOX)
    _, tf2 = rich_tf(s, Inches(7.1), Inches(1.75), Inches(5.0), Inches(2.3))
    add_para(tf2, "Narayanan (2022)", size=14, color=ACCENT3, bold=True, after=4)
    add_para(tf2, "The Limits of the Quantitative\nApproach to Discrimination",
             size=20, color=TITLE_CLR, bold=True, after=10)
    add_para(tf2,
             "\"Quantitative methods that treat observed outcomes as fair "
             "baselines risk laundering existing power structures into "
             "algorithmic objectivity.\"",
             size=14, color=BODY_CLR, italic=True, after=6)
    add_para(tf2, "Our laissez-faire baseline is not neutral \u2014 it reflects "
             "who currently has the resources to resist.",
             size=13, color=BODY_CLR, after=0)

    # Connection
    _, tf3 = rich_tf(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(2.5))
    add_para(tf3, "Our approach:", size=20, bold=True, color=TITLE_CLR, after=8)
    bullet(tf3,
           "We don\u2019t prescribe policy \u2014 we model the tradeoffs so "
           "the debate can be more informed", size=16, after=6)
    bullet(tf3,
           "By building both a data-driven model (XGBoost) and a theory-driven "
           "model (generative), we make visible the tension between quantitative "
           "prediction and qualitative reasoning about community consent",
           size=16, after=6)
    bullet(tf3,
           "The simulation renders distributional consequences visible: "
           "who gains, who loses, who pays", size=16, after=0)

    notes(s,
          "[Amrita]\n\n"
          "This project operates in two of the roles Abebe et al. describe. "
          "As a formalizer, we take the implicit power structure of data center "
          "siting and make it explicit and measurable. As a rebuttal, we simulate "
          "alternatives that challenge the assumption the status quo is neutral.\n\n"
          "Narayanan's caution is important here. Treating laissez-faire as an "
          "efficient baseline would launder existing power structures. Our simulation "
          "makes visible that the current system assigns all resistance costs to "
          "communities \u2014 that's not neutrality, it's a policy choice.\n\n"
          "We also address this through methodology: by building two parallel "
          "approaches \u2014 one data-driven (XGBoost) and one theory-driven "
          "(generative) \u2014 we make the modeling choices themselves transparent. "
          "The data-driven approach generalizes from observed outcomes but inherits "
          "their biases; the theory-driven approach is explicit about its assumptions "
          "but may miss real-world complexity.")


def slide_04(prs):
    """Intervention & Voting Thresholds."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Intervention & Voting Thresholds",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Scenario table
    rows, cols = 6, 4
    shape = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.5),
                               Inches(11), Inches(3.0))
    table = shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(3.3)

    hdrs = ["Scenario", "Threshold", "Who Bears Cost?", "Description"]
    data = [
        ["S1: Laissez-faire", "None", "Community (resistance)",
         "Current system \u2014 no formal consent"],
        ["S2: Majority (50%)", "50%", "Community (voting)",
         "Simple majority approval required"],
        ["S3: Supermajority (75%)", "75%", "Community (voting)",
         "Three-quarters must approve"],
        ["S4: Firm consent (50%)", "50%", "Firm (tax + jobs)",
         "Firm invests in benefits to earn consent"],
        ["S5: Firm consent (75%)", "75%", "Firm (tax + jobs)",
         "Same mechanism, higher bar"],
    ]

    for j, h in enumerate(hdrs):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BODY_CLR
                p.font.name = FONT
                if j == 2 and "Firm" in val:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW1 if i % 2 == 0 else TABLE_ROW2

    # Speculative inversion explanation
    _, tf = rich_tf(s, Inches(1.0), Inches(4.8), Inches(11), Inches(2.2))
    add_para(tf, "The speculative inversion:", size=18, bold=True,
             color=TITLE_CLR, after=8)
    bullet(tf,
           "S1\u2013S3: Community bears all costs \u2014 organizing, legal fees, "
           "sustained campaigns. The firm has no obligation to earn consent.", size=15, after=5)
    bullet(tf,
           "S4\u2013S5: The firm must actively invest in local tax relief and "
           "employment guarantees to push approval above the threshold. "
           "The cost shifts from community to corporation.", size=15, after=5)
    bullet(tf,
           "The intervention mechanism: firms solve an optimization \u2014 "
           "what\u2019s the cheapest mix of tax + employment investment to "
           "earn consent?", size=15, after=0)

    notes(s,
          "[Lauren]\n\n"
          "We model five consent regimes. The first three are variations of the "
          "status quo: no consent requirement, then progressively higher voting "
          "thresholds. In all three, the community bears the cost.\n\n"
          "The key innovation is scenarios 4 and 5. Here the legal default flips: "
          "the firm must obtain affirmative local consent BEFORE building. To do "
          "this, the firm can invest in two benefit channels: tax relief for local "
          "residents, and employment guarantees. These are modeled as nonlinear "
          "functions of county saturation \u2014 the first data center gets a big "
          "boost, but the 10th barely moves the needle.\n\n"
          "The firm solves a 2-variable linear program each time: what's the "
          "cheapest combination of tax and employment investment to clear the "
          "threshold? If no feasible solution exists, the project doesn't get built.")


def slide_05(prs):
    """Datasets."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Datasets",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    sources_left = [
        ("FracTracker Alliance (July 2025)",
         "1,380 facilities \u2192 232 counties, 108 labeled\n"
         "(81 approved, 27 blocked)"),
        ("Census ACS 5-Year (2022)",
         "Population, income, education, unemployment,\n"
         "agricultural employment \u2014 3,222 counties"),
        ("Census QWI (NAICS 5182)",
         "DC-specific employment & growth\n"
         "2,007 counties with data"),
        ("WRI Aqueduct 4.0",
         "Water stress via area-weighted spatial join\n"
         "3,144 counties"),
    ]

    sources_right = [
        ("MIT Election Lab (2024)",
         "Partisan lean (% Republican presidential)\n"
         "3,151 counties"),
        ("Good Jobs First + NCSL",
         "State-level incentive generosity scores\n"
         "51 states"),
        ("Opposition databases",
         "Bryce Rejection DB + DataCenterWatch\n"
         "47 documented opposition cases"),
    ]

    y = Inches(1.5)
    for name, desc in sources_left:
        box(s, Inches(0.8), y, Inches(5.5), Inches(1.0), LIGHT_BOX)
        _, tf = rich_tf(s, Inches(1.0), y + Inches(0.08), Inches(5.1), Inches(0.85))
        add_para(tf, name, size=13, bold=True, color=ACCENT, after=2)
        add_para(tf, desc, size=11, color=BODY_CLR, after=0)
        y += Inches(1.12)

    y = Inches(1.5)
    for name, desc in sources_right:
        box(s, Inches(6.8), y, Inches(5.5), Inches(1.0), LIGHT_BOX)
        _, tf = rich_tf(s, Inches(7.0), y + Inches(0.08), Inches(5.1), Inches(0.85))
        add_para(tf, name, size=13, bold=True, color=ACCENT, after=2)
        add_para(tf, desc, size=11, color=BODY_CLR, after=0)
        y += Inches(1.12)

    # Bottom note
    box(s, Inches(6.8), y, Inches(5.5), Inches(1.0), DARK_BOX)
    _, tfn = rich_tf(s, Inches(7.0), y + Inches(0.08), Inches(5.1), Inches(0.85))
    add_para(tfn, "Nothing is fabricated", size=13, bold=True, color=WHITE, after=2)
    add_para(tfn, "Every dataset has a reproducible fetch script.\n"
             "All code is open source on GitHub.", size=11,
             color=RGBColor(0xAA, 0xA7, 0xA0), after=0)

    notes(s,
          "[Lauren]\n\n"
          "All our data comes from real public sources. FracTracker gives us the "
          "facility locations. After filtering to hyperscale and mega campus "
          "facilities (>100 MW), we have 232 counties, of which 108 have clear "
          "approved or blocked outcomes for training the model.\n\n"
          "We enrich each county with features from six additional sources: "
          "Census ACS for demographics, Census QWI for DC-specific employment, "
          "WRI Aqueduct for water stress measured by spatial join, MIT Election "
          "Lab for partisan lean, Good Jobs First for state incentive policies, "
          "and two opposition databases for pushback history.\n\n"
          "Every dataset has a Python script that fetches it from the original "
          "API or public source. You can reproduce the entire data pipeline from "
          "scratch in about 15 minutes.")


def slide_06(prs):
    """Pipeline."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Pipeline: Data \u2192 Model \u2192 Simulation",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Three-stage pipeline boxes
    stages = [
        (Inches(0.5), "1. Feature Matrix",
         "7 data sources\n\u2192 15 features per county\n\u2192 3,153 counties",
         ACCENT),
        (Inches(4.5), "2. Approval Model",
         "XGBoost on 108 labeled counties\n"
         "\u2192 Calibration (3 anchor points)\n"
         "\u2192 P(approve) for all counties",
         ACCENT2),
        (Inches(8.5), "3. Monte Carlo Simulation",
         "120 months \u00d7 10,000 draws\n"
         "\u2192 5 scenarios\n"
         "\u2192 Growth, distribution, cost",
         ACCENT3),
    ]

    for x, title, desc, clr in stages:
        box(s, x, Inches(1.6), Inches(3.8), Inches(2.0), clr)
        add_textbox(s, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5),
                    title, font_size=18, color=WHITE, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(1.2),
                    desc, font_size=13, color=RGBColor(0xFF, 0xE8, 0xD0))

    # Arrows between stages
    for x in [Inches(4.3), Inches(8.3)]:
        add_textbox(s, x, Inches(2.2), Inches(0.3), Inches(0.5),
                    "\u2192", font_size=28, color=MUTED_CLR, alignment=PP_ALIGN.CENTER)

    # Detail below
    _, tf = rich_tf(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.2))
    add_para(tf, "15 structural features (XGBoost input):", size=16, bold=True,
             color=TITLE_CLR, after=6)
    add_para(tf,
             "avg_project_mw  \u00b7  hyperscaler_share  \u00b7  pushback_flag  \u00b7  "
             "state_incentive_score  \u00b7  dc_employment  \u00b7  dc_employment_growth  "
             "\u00b7  water_stress_decile  \u00b7  partisan_lean_r  \u00b7  population  "
             "\u00b7  population_density  \u00b7  median_household_income  \u00b7  "
             "unemployment_rate  \u00b7  pct_college_educated  \u00b7  "
             "ag_employment_share  \u00b7  electricity_price",
             size=12, color=MUTED_CLR, after=10)

    add_para(tf, "Key design decision:", size=16, bold=True, color=TITLE_CLR, after=4)
    add_para(tf,
             "Structural features only in XGBoost (demographics, environment, politics). "
             "Dynamic effects (saturation, intervention benefits) are modeled separately "
             "during simulation to avoid double-counting.",
             size=14, color=BODY_CLR, after=0)

    notes(s,
          "[Lauren]\n\n"
          "The pipeline has three stages. First, we assemble a feature matrix "
          "from 7 data sources \u2014 15 structural features per county covering "
          "demographics, economics, environment, and politics.\n\n"
          "Second, we train an XGBoost classifier on the 108 labeled counties, "
          "calibrate against three external anchor points, and extrapolate to "
          "all 3,153 U.S. counties.\n\n"
          "Third, we run the Monte Carlo simulation: 120 monthly time steps "
          "across 5 scenarios, each with 10,000 draws.\n\n"
          "An important architectural choice: we only put structural features "
          "in the XGBoost model. Dynamic effects like saturation and intervention "
          "benefits are handled by separate nonlinear functions during simulation. "
          "This prevents double-counting and makes the causal structure transparent.")


def slide_07(prs):
    """Two Approaches: XGBoost vs Generative."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Two Approaches to Modeling Approval",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Left: XGBoost (data-driven)
    box(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.8), LIGHT_BOX)
    _, tf1 = rich_tf(s, Inches(1.1), Inches(1.65), Inches(5.0), Inches(4.5))
    add_para(tf1, "XGBoost", size=24, bold=True, color=ACCENT, after=2)
    add_para(tf1, "Data-Driven", size=16, color=MUTED_CLR, after=10)

    add_para(tf1, "How it works:", size=15, bold=True, color=TITLE_CLR, after=4)
    bullet(tf1, "Train on 108 counties with known outcomes", size=14, after=3)
    bullet(tf1, "Learn nonlinear feature\u2013approval relationships", size=14, after=3)
    bullet(tf1, "Rank feature importance from data", size=14, after=3)
    bullet(tf1, "Extrapolate to all 3,153 counties", size=14, after=10)

    add_para(tf1, "Strength:", size=14, bold=True, color=ACCENT2, after=2)
    add_para(tf1, "Generalizes from observed outcomes.\nNo manual parameter assumptions.",
             size=13, color=BODY_CLR, after=8)
    add_para(tf1, "Limitation:", size=14, bold=True, color=ACCENT3, after=2)
    add_para(tf1, "Inherits biases in training data.\nOverconfident outside training distribution.",
             size=13, color=BODY_CLR, after=0)

    # Right: Generative (theory-driven)
    box(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.8), LIGHT_BOX)
    _, tf2 = rich_tf(s, Inches(7.3), Inches(1.65), Inches(5.0), Inches(4.5))
    add_para(tf2, "Generative Model", size=24, bold=True, color=ACCENT3, after=2)
    add_para(tf2, "Theory-Driven", size=16, color=MUTED_CLR, after=10)

    add_para(tf2, "How it works:", size=15, bold=True, color=TITLE_CLR, after=4)
    bullet(tf2, "Start from domain knowledge & survey data", size=14, after=3)
    bullet(tf2, "Assign explicit penalties/bonuses per feature", size=14, after=3)
    bullet(tf2, "Construct approval from first principles", size=14, after=3)
    bullet(tf2, "Transparent causal assumptions", size=14, after=10)

    add_para(tf2, "Strength:", size=14, bold=True, color=ACCENT2, after=2)
    add_para(tf2, "Every assumption is explicit and debatable.\n"
             "Interpretable causal reasoning.",
             size=13, color=BODY_CLR, after=8)
    add_para(tf2, "Limitation:", size=14, bold=True, color=ACCENT3, after=2)
    add_para(tf2, "Assumes penalty structure from theory.\n"
             "May miss real-world nonlinearities.",
             size=13, color=BODY_CLR, after=0)

    # Bottom: connection
    box(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), DARK_BOX)
    add_textbox(s, Inches(1.1), Inches(6.55), Inches(11), Inches(0.6),
                "Connection to qualitative vs quantitative: the data-driven model learns "
                "from what happened; the generative model encodes reasoning about why. "
                "Together they make our modeling choices visible.",
                font_size=14, color=RGBColor(0xCC, 0xC8, 0xC0))

    notes(s,
          "[Dhruv / transition slide]\n\n"
          "We built two parallel approaches to modeling community approval.\n\n"
          "The XGBoost approach is data-driven: train on 108 counties, let the "
          "algorithm learn which features matter, then extrapolate. Its strength "
          "is that it generalizes from real outcomes without manual assumptions. "
          "Its weakness is that it inherits the biases of who got built where.\n\n"
          "The generative approach is theory-driven: start from domain knowledge "
          "and survey evidence, then construct approval probability by assigning "
          "explicit penalties and bonuses to features. Its strength is "
          "transparency \u2014 every assumption is debatable. Its weakness is that "
          "those assumptions may be wrong.\n\n"
          "This duality connects to the broader tension between quantitative and "
          "qualitative methods that Narayanan raises. The data-driven model tells "
          "us what correlates with approval; the theory-driven model forces us to "
          "articulate why. Neither alone is sufficient.")


def slide_08(prs):
    """XGBoost Method (detail)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "XGBoost: Data-Driven Approval Model",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Left: training details
    _, tf = rich_tf(s, Inches(1.0), Inches(1.5), Inches(5.5), Inches(5.0))
    add_para(tf, "Training", size=18, bold=True, color=TITLE_CLR, after=8)
    bullet(tf, "108 labeled counties (81 approved, 27 blocked)", size=14, after=4)
    bullet(tf, "Binary classification: P(county approves)", size=14, after=4)
    bullet(tf, "5-fold stratified cross-validation", size=14, after=4)
    bullet(tf, "CV AUC: 0.679 \u00b1 0.08", size=14, bold=True, after=10)

    add_para(tf, "Calibration", size=18, bold=True, color=TITLE_CLR, after=8)
    bullet(tf, "Raw XGBoost scores \u2192 linear rescaling", size=14, after=4)
    bullet(tf, "3 anchor points from external evidence:", size=14, after=3)
    bullet(tf, "National median: 44% (Heatmap/Embold survey)", size=12, level=1, after=2)
    bullet(tf, "Loudoun Co., VA: 77.5% (JLARC, historical)", size=12, level=1, after=2)
    bullet(tf, "Prince William, VA: 25% (moratorium)", size=12, level=1, after=6)
    bullet(tf, "State-level shrinkage (K=5) for counties outside "
           "training distribution", size=14, after=0)

    # Right: feature importance
    box(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), LIGHT_BOX)
    _, tf2 = rich_tf(s, Inches(7.3), Inches(1.65), Inches(5.0), Inches(4.7))
    add_para(tf2, "Feature Importance (top 10)", size=18, bold=True,
             color=TITLE_CLR, after=8)

    features = [
        ("pushback_flag", "0.310", "Direct opposition signal"),
        ("hyperscaler_share", "0.120", "Big Tech vs colo operator"),
        ("state_incentive_score", "0.098", "Policy environment"),
        ("dc_employment", "0.091", "Existing DC workforce"),
        ("avg_project_mw", "0.076", "Project size/visibility"),
        ("partisan_lean_r", "0.065", "Political baseline"),
        ("water_stress_decile", "0.058", "Environmental strain"),
        ("pct_college_educated", "0.048", "Opposition capacity"),
        ("population_density", "0.042", "Urban/rural character"),
        ("unemployment_rate", "0.038", "Receptivity to jobs"),
    ]

    for fname, imp, desc in features:
        add_para(tf2, f"{fname}: {imp}", size=12, bold=True, color=ACCENT, after=0)
        add_para(tf2, f"  {desc}", size=10, color=MUTED_CLR, after=4)

    notes(s,
          "[Dhruv]\n\n"
          "Here's the XGBoost model in detail. We train on 108 counties: "
          "81 that approved data centers and 27 that blocked them.\n\n"
          "The model achieves a cross-validation AUC of 0.679. With 108 samples "
          "and 15 features, this is a reasonable result \u2014 not great, but "
          "meaningful. Small changes in regularization dramatically affect "
          "performance at this sample size, so we kept things simple.\n\n"
          "After training, we calibrate the raw scores against three external "
          "anchor points: the national median at 44% from the Heatmap/Embold "
          "survey, Loudoun County at 77.5% from JLARC data, and Prince William "
          "at 25% reflecting its moratorium.\n\n"
          "Feature importance tells an interesting story. Pushback flag dominates "
          "at 0.31 \u2014 counties with documented opposition are very different "
          "from those without. Hyperscaler share matters: counties with Big Tech "
          "operators behave differently from those with colocation providers. "
          "State incentive policy is third \u2014 the policy environment matters.\n\n"
          "We also apply state-level shrinkage: counties in states with few "
          "training examples are pulled toward the 44% national median. This "
          "prevents overconfident extrapolation to the 2,900+ counties outside "
          "the training distribution.")


def slide_09(prs):
    """USA Map \u2014 County Approval Probability."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "County Approval Probability Map",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Map placeholder
    img_placeholder(s, Inches(0.8), Inches(1.4), Inches(9.0), Inches(5.0),
                    "screenshot of full_approval_map.html")

    # Callout boxes
    y = Inches(1.4)
    callouts = [
        ("Loudoun Co., VA", "77.5%", ACCENT2),
        ("National median", "44%", ACCENT),
        ("Prince William, VA", "25%", ACCENT3),
    ]
    for label, pct, clr in callouts:
        box(s, Inches(10.2), y, Inches(2.5), Inches(1.2), LIGHT_BOX)
        _, tf = rich_tf(s, Inches(10.35), y + Inches(0.1), Inches(2.2), Inches(1.0))
        add_para(tf, pct, size=28, bold=True, color=clr, after=2, align=PP_ALIGN.CENTER)
        add_para(tf, label, size=12, color=MUTED_CLR, after=0, align=PP_ALIGN.CENTER)
        y += Inches(1.45)

    add_textbox(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
                "XGBoost trained on 108 counties \u2192 calibrated with 3 anchors \u2192 "
                "extrapolated to 3,153 counties (with state-level shrinkage for OOS counties)",
                font_size=13, color=MUTED_CLR)

    notes(s,
          "[Dhruv]\n\n"
          "This is the XGBoost model's predicted approval probability for every "
          "U.S. county. Green = likely to approve, red = likely to block.\n\n"
          "The established data center corridors \u2014 Virginia, Georgia, Texas \u2014 "
          "show the highest approval. The Midwest and rural areas are closer to "
          "the national median.\n\n"
          "Caveat: for the roughly 2,900 greenfield counties with no facility "
          "history, predictions are driven by demographics, water stress, and "
          "partisan lean alone. We apply state-level shrinkage to be conservative "
          "\u2014 counties in under-represented states are pulled toward the 44% "
          "median rather than being given extreme predictions.")


def slide_10(prs):
    """Generative Approval Model (done outside this project)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Generative Approval Model",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Placeholder structure
    box(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.2), LIGHT_BOX)
    _, tf_intro = rich_tf(s, Inches(1.3), Inches(1.6), Inches(10.7), Inches(1.0))
    add_para(tf_intro, "Theory-driven approach: construct approval probability from "
             "first principles using domain knowledge and survey evidence, rather than "
             "learning from observed outcomes.",
             size=16, color=BODY_CLR, after=4)
    add_para(tf_intro, "(Developed separately from the XGBoost pipeline)",
             size=13, color=MUTED_CLR, italic=True, after=0)

    # Content area - placeholder for teammate to fill
    box(s, Inches(1.0), Inches(3.0), Inches(5.5), Inches(3.5), LIGHT_BOX)
    _, tf_left = rich_tf(s, Inches(1.3), Inches(3.15), Inches(5.0), Inches(3.2))
    add_para(tf_left, "Approach", size=18, bold=True, color=ACCENT3, after=6)
    bullet(tf_left, "[How baseline approval is constructed]", size=14,
           color=MUTED_CLR, after=4)
    bullet(tf_left, "[Feature penalties/bonuses assumed]", size=14,
           color=MUTED_CLR, after=4)
    bullet(tf_left, "[Calibration to survey evidence]", size=14,
           color=MUTED_CLR, after=4)
    bullet(tf_left, "[Coverage: which counties?]", size=14,
           color=MUTED_CLR, after=4)

    box(s, Inches(7.0), Inches(3.0), Inches(5.3), Inches(3.5), LIGHT_BOX)
    _, tf_right = rich_tf(s, Inches(7.3), Inches(3.15), Inches(4.8), Inches(3.2))
    add_para(tf_right, "Comparison with XGBoost", size=18, bold=True,
             color=ACCENT3, after=6)
    bullet(tf_right, "[Where the two models agree/disagree]", size=14,
           color=MUTED_CLR, after=4)
    bullet(tf_right, "[What the generative model makes explicit\n"
           "that XGBoost hides in learned weights]", size=14,
           color=MUTED_CLR, after=4)
    bullet(tf_right, "[Connection to Narayanan: making\nassumptions visible]",
           size=14, color=MUTED_CLR, after=4)

    notes(s,
          "[Teammate who built the generative model]\n\n"
          "FILL IN: Explain the theory-driven generative approval model.\n\n"
          "Key points to cover:\n"
          "- How is baseline approval probability constructed?\n"
          "- What penalties/bonuses are assumed for each feature?\n"
          "- How does this connect to the qualitative vs quantitative tension?\n"
          "- Where does it agree with and diverge from the XGBoost predictions?\n"
          "- What does it make visible that the data-driven approach hides?")


def slide_11(prs):
    """Monte Carlo Simulation Engine."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Monte Carlo Simulation Engine",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Left: simulation steps
    steps = [
        ("1", "Generate candidates",
         "~5 projects/month at 300 MW each\n(1.5 GW/month national rate)"),
        ("2", "Assign to counties",
         "State share \u00d7 county weight\n(existing facility counties 3\u00d7 more likely)"),
        ("3", "Sample approval",
         "Beta(\u03b1, \u03b2) centered on county's\ncalibrated probability (\u03ba = 40)"),
        ("4", "Apply regime rule",
         "Laissez-faire: draw > 0.5\nVoting: draw > threshold\n"
         "Firm-borne: LP-optimized investment"),
        ("5", "Update saturation",
         "If built: n \u2192 n+1 for that county\n"
         "Intervention deltas decay with n"),
    ]

    y = Inches(1.5)
    for num, title, desc in steps:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05),
                                  Inches(0.4), Inches(0.4))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        add_textbox(s, Inches(0.8), y + Inches(0.05), Inches(0.4), Inches(0.4),
                    num, font_size=14, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        _, tf = rich_tf(s, Inches(1.4), y, Inches(4.8), Inches(0.95))
        add_para(tf, title, size=15, bold=True, color=TITLE_CLR, after=2)
        add_para(tf, desc, size=11, color=MUTED_CLR, after=0)
        y += Inches(1.0)

    # Right: intervention functions
    box(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3.5), LIGHT_BOX)
    _, tfi = rich_tf(s, Inches(7.1), Inches(1.65), Inches(5.0), Inches(3.2))
    add_para(tfi, "Intervention Functions (S4 & S5 only)", size=16, bold=True,
             color=TITLE_CLR, after=8)
    add_para(tfi, "Tax relief (exponential decay):", size=13, bold=True,
             color=ACCENT, after=2)
    add_para(tfi, "\u0394p = 0.20 \u00d7 e^(\u22120.25n)\n"
             "Strong early, fades with saturation.\n"
             "$405.8M annual tax revenue per GW.", size=11, color=BODY_CLR, after=8)
    add_para(tfi, "Employment (bell curve):", size=13, bold=True,
             color=ACCENT2, after=2)
    add_para(tfi, "\u0394p = 0.15 \u00d7 (n/10) \u00d7 e^(1\u2212n/10)\n"
             "Peaks at ~10 facilities, then fades.\n"
             "45K construction + 5K permanent jobs per GW.", size=11,
             color=BODY_CLR, after=6)

    img_placeholder(s, Inches(6.8), Inches(5.2), Inches(5.5), Inches(1.8),
                    "Figure 2: intervention curves from proposal")

    # Scale
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(5.5), Inches(0.5),
                "120 months  \u00d7  10,000 draws  \u00d7  5 scenarios",
                font_size=16, color=MUTED_CLR, bold=True)

    notes(s,
          "[Dhruv]\n\n"
          "Here's how the simulation works step by step.\n\n"
          "Each month, about 5 candidate projects enter the pipeline. They're "
          "distributed to states proportional to current DC electricity share "
          "from EIA data, then to counties weighted 3\u00d7 toward places with "
          "existing facilities.\n\n"
          "For each candidate, we sample an approval draw from a Beta distribution "
          "centered on that county's probability. The concentration parameter "
          "kappa=40 means the draw typically stays within \u00b110% of the mean.\n\n"
          "Then the scenario rule applies. In firm-borne scenarios, the firm "
          "solves a linear program: minimize the cost of tax relief + employment "
          "investment subject to pushing expected approval above the threshold.\n\n"
          "The intervention curves are the key mechanism. Tax benefits follow "
          "exponential decay \u2014 the first DC gets a 20% approval boost, the 10th "
          "gets almost nothing. Employment benefits follow a bell curve peaking "
          "around 10 facilities. Both are calibrated to real figures from CEA "
          "Florida and JLARC Virginia.\n\n"
          "Geographic substitution: if a project is blocked, there's a 50% chance "
          "the firm tries another county in the same state.")


def slide_12(prs):
    """Current Snapshot: What Vanishes? (generative)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Current Snapshot: What Vanishes?",
                font_size=36, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Explanation
    box(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.0), LIGHT_BOX)
    _, tf_exp = rich_tf(s, Inches(1.3), Inches(1.5), Inches(10.7), Inches(0.8))
    add_para(tf_exp,
             "Under each consent regime, which counties lose feasibility? "
             "Using the generative model\u2019s transparent assumptions, "
             "we can trace exactly which features cause a county to fall "
             "below threshold.",
             size=15, color=BODY_CLR, after=0)

    # Main content area - placeholder
    img_placeholder(s, Inches(1.0), Inches(2.7), Inches(7.0), Inches(3.8),
                    "generative model: map or chart showing\n"
                    "which counties lose feasibility under each threshold")

    # Right: key insights placeholder
    box(s, Inches(8.5), Inches(2.7), Inches(4.0), Inches(3.8), LIGHT_BOX)
    _, tf_r = rich_tf(s, Inches(8.7), Inches(2.85), Inches(3.6), Inches(3.5))
    add_para(tf_r, "Key Insights", size=16, bold=True, color=ACCENT3, after=8)
    bullet(tf_r, "[Which feature penalties cause the\nmost counties to drop out?]",
           size=13, color=MUTED_CLR, after=6)
    bullet(tf_r, "[How does the generative model\u2019s\n"
           "\"what vanishes\" compare to\nXGBoost predictions?]",
           size=13, color=MUTED_CLR, after=6)
    bullet(tf_r, "[What does this reveal about\nthe cost of consent regimes\n"
           "that isn\u2019t visible in aggregate\nGW numbers?]",
           size=13, color=MUTED_CLR, after=0)

    notes(s,
          "[Teammate who built the generative model]\n\n"
          "FILL IN: Using the generative model, show what happens to specific "
          "counties under each consent regime.\n\n"
          "Key questions to answer:\n"
          "- Which counties are feasible under laissez-faire but vanish under "
          "a 50% threshold? A 75% threshold?\n"
          "- What features drive counties below threshold?\n"
          "- How does this compare to the XGBoost predictions?\n"
          "- What does this reveal about the distributional impact of consent "
          "regimes that isn't visible in aggregate numbers?")


def slide_13(prs):
    """Projection: Future Growth Map (XGBoost)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Projection: Where Data Centers Get Built (2026\u20132035)",
                font_size=34, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # Two maps side by side
    img_placeholder(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.0),
                    "s1_evolution.gif\n(Laissez-faire)")
    add_textbox(s, Inches(0.5), Inches(5.5), Inches(6.0), Inches(0.4),
                "S1: Laissez-faire  |  66.3 GW  |  Gini = 0.956",
                font_size=13, color=BODY_CLR, bold=True, alignment=PP_ALIGN.CENTER)

    img_placeholder(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(4.0),
                    "s4_evolution.gif\n(Firm consent 50%)")
    add_textbox(s, Inches(6.8), Inches(5.5), Inches(6.0), Inches(0.4),
                "S4: Firm consent 50%  |  107.5 GW  |  Gini = 0.925",
                font_size=13, color=BODY_CLR, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom observation
    _, tf = rich_tf(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(1.2))
    add_para(tf, "Geographic concentration persists in all scenarios (Gini > 0.92). "
             "Firm investment increases total capacity and spreads it marginally, "
             "but doesn\u2019t overcome the structural advantages of established "
             "data center corridors (Virginia, Georgia, Texas).",
             size=14, color=MUTED_CLR, after=0)

    notes(s,
          "[Dhruv]\n\n"
          "These animated maps show the XGBoost simulation results \u2014 where "
          "data centers actually get built over 120 months. Circle size grows "
          "as more facilities accumulate in a county.\n\n"
          "Under laissez-faire, building concentrates heavily in Virginia, "
          "Georgia, and Texas. The Gini coefficient is 0.956 \u2014 extreme "
          "concentration.\n\n"
          "Under firm consent at 50%, you see more circles in more places. "
          "But I want to be honest: Gini only drops to 0.925. Still extreme "
          "concentration. Firm investment marginally improves distribution "
          "but doesn't overcome structural advantages.\n\n"
          "What DOES change is total volume: 107.5 GW vs 66.3 GW. And every "
          "build in S4 comes with contractual community benefits.")


def slide_14(prs):
    """Projection: Future Growth Curves (XGBoost)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_textbox(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.7),
                "Projection: Cumulative Growth Trajectories",
                font_size=34, color=TITLE_CLR, bold=True)
    accent_line(s, Inches(1.0), Inches(1.15), Inches(1.5))

    # GIF placeholder
    img_placeholder(s, Inches(0.8), Inches(1.4), Inches(7.5), Inches(4.8),
                    "gw_growth_comparison.gif")

    # Key findings sidebar
    y = Inches(1.4)

    box(s, Inches(8.8), y, Inches(4.0), Inches(1.3), ACCENT)
    _, tf1 = rich_tf(s, Inches(9.0), y + Inches(0.1), Inches(3.6), Inches(1.1))
    add_para(tf1, "S1 = S2: Identical", size=15, bold=True, color=WHITE, after=3)
    add_para(tf1, "Formalizing 50% changes nothing.\nThe system already has an\n"
             "implicit majority threshold.",
             size=11, color=RGBColor(0xFF, 0xE8, 0xD0), after=0)
    y += Inches(1.5)

    box(s, Inches(8.8), y, Inches(4.0), Inches(1.1), LIGHT_BOX)
    _, tf2 = rich_tf(s, Inches(9.0), y + Inches(0.1), Inches(3.6), Inches(0.9))
    add_para(tf2, "S4: 107.5 GW (+62%)", size=15, bold=True, color=ACCENT2, after=2)
    add_para(tf2, "Firm investment unlocks growth\nbeyond the baseline.",
             size=11, color=BODY_CLR, after=0)
    y += Inches(1.3)

    box(s, Inches(8.8), y, Inches(4.0), Inches(1.1), LIGHT_BOX)
    _, tf3 = rich_tf(s, Inches(9.0), y + Inches(0.1), Inches(3.6), Inches(0.9))
    add_para(tf3, "S3: 9.0 GW (\u221286%)", size=15, bold=True, color=ACCENT3, after=2)
    add_para(tf3, "Supermajority without investment\nnearly kills all development.",
             size=11, color=BODY_CLR, after=0)
    y += Inches(1.3)

    box(s, Inches(8.8), y, Inches(4.0), Inches(1.1), LIGHT_BOX)
    _, tf4 = rich_tf(s, Inches(9.0), y + Inches(0.1), Inches(3.6), Inches(0.9))
    add_para(tf4, "S5: 43.5 GW", size=15, bold=True, color=BODY_CLR, after=2)
    add_para(tf4, "High bar + firm investment:\nless capacity, but consent is earned.",
             size=11, color=BODY_CLR, after=0)

    # Results table (compact)
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
                "S4 firm cost: $11.4B  |  S5 firm cost: $10.2B  "
                "\u2014  nearly the same cost, 60% less capacity at the higher threshold",
                font_size=13, color=MUTED_CLR)

    notes(s,
          "[Dhruv]\n\n"
          "This is the headline result slide. The growth curves show cumulative "
          "U.S. data center capacity over 10 years under each scenario, with "
          "95% confidence intervals from 10,000 Monte Carlo draws.\n\n"
          "Four key findings:\n\n"
          "First: S1 and S2 produce IDENTICAL results. 66.3 GW, 221 facilities. "
          "The laissez-faire scenario already uses approval_draw > 0.5. Adding a "
          "formal 50% vote changes literally nothing. This is perhaps our most "
          "telling finding \u2014 formalizing what already exists is not an "
          "intervention.\n\n"
          "Second: S4 (firm consent 50%) INCREASES capacity by 62%. When firms "
          "invest in benefits, they enable the intervention functions that boost "
          "approval above 50% in counties that would otherwise block.\n\n"
          "Third: S3 (supermajority 75%, no investment) nearly kills development. "
          "Only 9 GW. Almost no county has baseline approval above 75%.\n\n"
          "Fourth: S5 is a middle ground. 43.5 GW, but 63% of candidates are "
          "infeasible \u2014 no amount of investment can push most counties to 75%.\n\n"
          "The firms spend almost the same total ($11.4B vs $10.2B) regardless "
          "of threshold, but the higher bar produces 60% less capacity. There's "
          "a sweet spot.")


def slide_15(prs):
    """Conclusion."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BOX)

    accent_line(s, Inches(1.5), Inches(0.8), Inches(2))

    add_textbox(s, Inches(1.5), Inches(1.1), Inches(10), Inches(0.8),
                "Conclusion",
                font_size=38, color=WHITE, bold=True)

    findings = [
        ("1", "Consent doesn\u2019t have to kill growth",
         "Firm-borne consent at 50% increases capacity by 62%. When firms invest "
         "in communities, they unlock counties that would otherwise block development.",
         ACCENT),
        ("2", "The threshold matters enormously",
         "50% vs 75% = 107.5 GW vs 43.5 GW. And formalizing 50% without firm "
         "investment (S2) changes literally nothing \u2014 the system already has "
         "an implicit majority threshold.",
         ACCENT2),
        ("3", "The current system is not neutral",
         "Laissez-faire assigns all resistance costs to communities. Our simulation "
         "makes this visible \u2014 not to prescribe policy, but so the tradeoffs "
         "can be debated honestly.",
         ACCENT3),
    ]

    y = Inches(2.3)
    for num, title, desc, clr in findings:
        add_textbox(s, Inches(1.5), y, Inches(0.5), Inches(0.5),
                    num, font_size=28, color=clr, bold=True)
        _, tf = rich_tf(s, Inches(2.2), y, Inches(9.5), Inches(1.3))
        add_para(tf, title, size=20, bold=True, color=WHITE, after=4)
        add_para(tf, desc, size=14, color=RGBColor(0xAA, 0xA7, 0xA0), after=0)
        y += Inches(1.5)

    # Final line
    add_textbox(s, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
                "\"What if\" is not a prediction \u2014 it\u2019s an invitation "
                "to imagine a different default.",
                font_size=16, color=RGBColor(0x88, 0x85, 0x80), italic=True)

    notes(s,
          "[Tyler / everyone]\n\n"
          "Three takeaways:\n\n"
          "First: consent and growth aren't opposed. The counterintuitive finding "
          "is that moderate consent requirements paired with firm investment produce "
          "MORE total capacity than the status quo. The 'do nothing' approach isn't "
          "efficient \u2014 it just hides the costs.\n\n"
          "Second: the threshold is the most important policy lever. 50% is a sweet "
          "spot. And the S1=S2 identity is telling: adding a formal vote without "
          "changing the cost structure accomplishes nothing.\n\n"
          "Third: this is the Narayanan point and the Abebe point combined. The "
          "current system isn't neutral. By modeling it explicitly alongside "
          "alternatives, we make the distributional consequences visible. We're "
          "not prescribing firm-borne consent as policy \u2014 we're showing what "
          "the tradeoffs look like when the cost burden shifts.\n\n"
          "This is speculative algorithm design: using computation not to optimize, "
          "but to ask 'what if?'")


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    builders = [
        slide_01, slide_02, slide_03, slide_04, slide_05,
        slide_06, slide_07, slide_08, slide_09, slide_10,
        slide_11, slide_12, slide_13, slide_14, slide_15,
    ]
    for fn in builders:
        fn(prs)

    out = Path(__file__).resolve().parent.parent / "outputs" / "presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"  {len(prs.slides)} slides with speaker notes")
    print(f"\nImport into Google Slides:")
    print(f"  File > Open > Upload > presentation.pptx")
    print(f"\nPlaceholders to fill:")
    print(f"  Slide 9:  Screenshot of full_approval_map.html")
    print(f"  Slide 10: Generative model content (teammate)")
    print(f"  Slide 11: Figure 2 (intervention curves)")
    print(f"  Slide 12: Generative model 'what vanishes' (teammate)")
    print(f"  Slide 13: s1_evolution.gif + s4_evolution.gif")
    print(f"  Slide 14: gw_growth_comparison.gif")


if __name__ == "__main__":
    main()
